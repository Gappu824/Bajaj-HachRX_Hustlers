# app/core/rag_pipeline.py - Optimized version
import io
import re
import os
import logging
import requests
import asyncio
import aiohttp
import time
import json
import hashlib
import msgpack
import lz4.frame
from typing import List, Tuple, Dict, Set, Optional, Any
from collections import defaultdict
import zipfile

# Document parsing imports
import pdfplumber
from docx import Document
from odf.text import P
from odf.opendocument import load
import PyPDF2
import pypdf
import pandas as pd
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available, PowerPoint parsing disabled")
try:
    from PIL import Image
    import pytesseract
    IMAGE_PROCESSING_AVAILABLE = True
except ImportError:
    IMAGE_PROCESSING_AVAILABLE = False
    logging.warning("Pillow or pytesseract not available, image processing disabled")    

# ML and AI imports
from sentence_transformers import SentenceTransformer
from sentence_transformers.cross_encoder import CrossEncoder
import faiss
import numpy as np
import google.generativeai as genai
from tenacity import retry, stop_after_attempt, wait_exponential, retry_if_exception_type
from google.api_core import exceptions as google_exceptions

# Local imports
from app.core.config import settings
from app.core.cache import cache
from app.core.enhanced_retrieval import EnhancedRetriever
from app.core.contextual_embeddings import ContextualRetrieval
from concurrent.futures import ThreadPoolExecutor, ProcessPoolExecutor
import multiprocessing

logger = logging.getLogger(__name__)


class SmartChunker:
    """Intelligent chunking that preserves semantic boundaries"""

    # In SmartChunker class, add this method:
    @staticmethod
    def chunk_spreadsheet_data(text: str, metadata: List[Dict]) -> Tuple[List[str], List[Dict]]:
        """Special chunking for spreadsheet data to preserve row context"""
        chunks = []
        chunk_metadata = []
        
        lines = text.split('\n')
        
        # Keep headers in every chunk for context
        headers = []
        summary = []
        data_lines = []
        
        for line in lines:
            if line.startswith('Column Headers:') or line.startswith('Headers:'):
                headers.append(line)
            elif line.startswith('===') or line.startswith('Shape:') or line.startswith('Total'):
                summary.append(line)
            elif line.startswith('Row'):
                data_lines.append(line)
        
        # Create chunks with headers included
        header_text = '\n'.join(headers + summary) if headers else ''
        
        # Chunk data rows
        chunk_size = 20  # rows per chunk
        for i in range(0, len(data_lines), chunk_size):
            chunk_lines = data_lines[i:i+chunk_size]
            if header_text:
                chunk_text = header_text + '\n\n' + '\n'.join(chunk_lines)
            else:
                chunk_text = '\n'.join(chunk_lines)
            
            chunks.append(chunk_text)
            chunk_metadata.append({
                'type': 'spreadsheet_chunk',
                'start_row': i + 1,
                'end_row': min(i + chunk_size, len(data_lines))
            })
        
        return chunks, chunk_metadata if chunks else ([text], metadata)
    
    @staticmethod
    def chunk_with_boundaries(text: str, metadata: List[Dict], 
                            chunk_size: int = 800,
                            overlap: int = 100) -> Tuple[List[str], List[Dict]]:
        """Create chunks with adaptive sizing based on document length"""
        
        # Adaptive chunk size based on document length
        doc_length = len(text)
        if doc_length < 5000:
            chunk_size = 500
            overlap = 50
        elif doc_length < 20000:
            chunk_size = 800
            overlap = 100
        else:
            chunk_size = 1200
            overlap = 150
        
        chunks = []
        chunk_metadata = []
        
        # Fast path for very small documents
        if doc_length < 1000:
            chunks.append(text)
            chunk_metadata.append({'type': 'full_document'})
            return chunks, chunk_metadata
        
        # Split by major sections
        sections = re.split(r'\n(?=(?:CHAPTER|SECTION|ARTICLE|PART)\s+[IVXLCDM\d]+)', text)
        
        for section in sections:
            if not section.strip():
                continue
            
            # Split by paragraphs
            paragraphs = re.split(r'\n\n+', section)
            current_chunk = ""
            current_meta = {}
            
            for para in paragraphs:
                para = para.strip()
                if not para:
                    continue
                
                # Check if adding this paragraph exceeds chunk size
                if len(current_chunk) + len(para) + 2 > chunk_size and current_chunk:
                    chunks.append(current_chunk.strip())
                    chunk_metadata.append(current_meta)
                    
                    # Start new chunk with overlap
                    if len(current_chunk) > overlap:
                        overlap_text = current_chunk[-overlap:]
                        current_chunk = overlap_text + "\n\n" + para
                    else:
                        current_chunk = para
                    current_meta = SmartChunker._extract_metadata(para)
                else:
                    if current_chunk:
                        current_chunk += "\n\n" + para
                    else:
                        current_chunk = para
                        current_meta = SmartChunker._extract_metadata(para)
            
            # Add final chunk
            if current_chunk.strip():
                chunks.append(current_chunk.strip())
                chunk_metadata.append(current_meta)
        
        return chunks, chunk_metadata
    
    @staticmethod
    def _extract_metadata(text: str) -> Dict:
        """Extract metadata from text"""
        metadata = {'type': 'text'}
        
        # Check for headers
        if re.match(r'^(?:CHAPTER|SECTION|ARTICLE|PART)\s+[IVXLCDM\d]+', text):
            metadata['type'] = 'header'
            metadata['level'] = 1
        elif re.match(r'^\d+\.\s+[A-Z]', text):
            metadata['type'] = 'section'
            metadata['level'] = 2
        
        # Check for lists
        if re.match(r'^(?:\d+\.|[a-z]\)|•|-|\*)\s+', text):
            metadata['type'] = 'list_item'
        
        # Check for definitions
        if 'means' in text or 'defined as' in text or 'refers to' in text:
            metadata['type'] = 'definition'
        
        return metadata


class OptimizedVectorStore:
    """Enhanced vector store with hybrid retrieval"""
    
    def __init__(self, chunks: List[str], embeddings: np.ndarray, 
                 model: SentenceTransformer, chunk_metadata: Optional[List[Dict]] = None):
        self.chunks = chunks
        self.model = model
        self.chunk_metadata = chunk_metadata or [{}] * len(chunks)
        self.embeddings = embeddings
        dimension = embeddings.shape[1]
        
        # Use simple flat index for speed (good for < 100k chunks)
        self.index = faiss.IndexFlatL2(dimension)
        self.index.add(embeddings)
        
        # Initialize enhanced retriever only if we have multiple chunks
        if len(chunks) > 1:
            self.enhanced_retriever = EnhancedRetriever(chunks, chunk_metadata)
        else:
            self.enhanced_retriever = None
        
        # Build question patterns
        self.question_patterns = self._build_question_patterns()
        
        logger.info(f"Built optimized vector store with {len(chunks)} chunks")
    
    def _build_question_patterns(self) -> Dict[str, Dict]:
        """Build patterns for question classification"""
        patterns = {
            'definitional': {
                'patterns': [
                    r'what is', r'define', r'meaning of', r'explain the term',
                    r'definition of', r'what does .* mean'
                ],
                'boost_terms': ['is defined as', 'refers to', 'means', 'definition', 'is called']
            },
            'computational': {
                'patterns': [
                    r'calculate', r'how much', r'what is the total', r'compute',
                    r'what.*cost', r'what.*price', r'what.*amount'
                ],
                'boost_terms': ['formula', 'calculation', 'total', 'sum', '%', 'percentage', 'equals']
            },
            'procedural': {
                'patterns': [
                    r'how to', r'what is the process', r'steps to', r'explain the procedure',
                    r'what.*procedure', r'how do'
                ],
                'boost_terms': ['process', 'procedure', 'steps', 'first', 'then', 'finally', 'instructions']
            },
            'list_based': {
                'patterns': [
                    r'list all', r'what are all', r'mention the', r'enumerate',
                    r'what are the.*types', r'give.*examples'
                ],
                'boost_terms': ['include:', 'are as follows', 'following', 'such as', 'for example', 'includes']
            },
            'conditional': {
                'patterns': [
                    r'what if', r'under what circumstances', r'conditions',
                    r'when.*covered', r'is.*covered'
                ],
                'boost_terms': ['if', 'subject to', 'provided that', 'conditions', 'unless', 'except']
            }
        }
        return patterns
    
    def hybrid_search(self, query: str, k: int = 15) -> List[Tuple[str, float, Dict]]:
        """Fast hybrid search combining multiple retrieval methods"""
        # Detect question type
        question_type = self._detect_question_type(query.lower())
        
        # For single chunk documents, return immediately
        if len(self.chunks) == 1:
            return [(self.chunks[0], 1.0, {
                'chunk_idx': 0,
                'metadata': self.chunk_metadata[0],
                'question_type': question_type,
                'retrieval_method': 'single_chunk'
            })]
        
        # Get candidates from enhanced retriever if available
        if self.enhanced_retriever:
            bm25_candidates = self.enhanced_retriever.retrieve(
                query, k=min(k*2, len(self.chunks)), question_type=question_type
            )
        else:
            bm25_candidates = []
        
        # Get candidates from vector search
        query_embedding = self.model.encode([query]).astype('float32')
        search_k = min(k*2, len(self.chunks))
        distances, indices = self.index.search(query_embedding, search_k)
        
        # Combine and deduplicate results
        candidate_scores = defaultdict(float)
        
        # Add BM25 scores
        if bm25_candidates:
            max_bm25_score = max(score for _, score in bm25_candidates) if bm25_candidates else 1.0
            for idx, score in bm25_candidates:
                if idx < len(self.chunks):
                    candidate_scores[idx] += (score / max_bm25_score) * 0.5
        
        # Add vector similarity scores
        for dist, idx in zip(distances[0], indices[0]):
            if idx != -1 and idx < len(self.chunks):
                similarity = 1 / (1 + dist)
                candidate_scores[idx] += similarity * 0.5
        
        # Apply question-type specific boosting
        for idx in candidate_scores:
            chunk = self.chunks[idx]
            chunk_lower = chunk.lower()
            
            if question_type in self.question_patterns:
                boost_terms = self.question_patterns[question_type]['boost_terms']
                boost_score = sum(0.1 for term in boost_terms if term in chunk_lower)
                candidate_scores[idx] += boost_score
        
        # Sort by combined score
        sorted_candidates = sorted(candidate_scores.items(), key=lambda x: x[1], reverse=True)
        
        # Return top-k results
        results = []
        for idx, score in sorted_candidates[:k]:
            results.append((
                self.chunks[idx],
                score,
                {
                    'chunk_idx': idx,
                    'metadata': self.chunk_metadata[idx],
                    'question_type': question_type,
                    'retrieval_method': 'hybrid'
                }
            ))
        
        return results
    
    def _detect_question_type(self, query: str) -> str:
        """Detect question type"""
        for q_type, config in self.question_patterns.items():
            for pattern in config['patterns']:
                if re.search(pattern, query, re.IGNORECASE):
                    return q_type
        return 'general'


class HybridFastTrackRAGPipeline:
    """Optimized pipeline for speed and accuracy"""
    
    def __init__(self, embedding_model: SentenceTransformer):
        self.embedding_model = embedding_model
        self.cross_encoder = None
        
        # Configure Gemini
        try:
            genai.configure(api_key=settings.GOOGLE_API_KEY)
            logger.info("Google Generative AI configured successfully")
        except Exception as e:
            logger.critical(f"Failed to configure Google AI: {e}")
            raise
        
        # Thread pools
        self.thread_pool = ThreadPoolExecutor(max_workers=8)

    # In the HybridFastTrackRAGPipeline class within app/core/rag_pipeline.py

    # In the HybridFastTrackRAGPipeline class within app/core/rag_pipeline.py

    async def get_text_and_metadata(self, content: bytes, url: str) -> Tuple[str, List[Dict]]:
        """Helper function to parse content based on file extension from a URL."""
        file_extension = os.path.splitext(url.split('?')[0])[1].lower()
        
        if not file_extension or file_extension == '.pdf':
            return self._parse_pdf_enhanced(content)
        elif file_extension in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff']:
            return self._parse_image(content)
        # This function now handles all other formats.
        else:
            return self._parse_other_formats(content, file_extension)

    async def _parse_zip(self, content: bytes) -> Tuple[str, List[Dict]]:
        """
        Extracts files from a ZIP archive and processes their content recursively.
        """
        logger.info("Processing ZIP archive...")
        combined_text = []
        combined_metadata = []

        with io.BytesIO(content) as temp_file:
            with zipfile.ZipFile(temp_file, 'r') as zip_ref:
                file_list = zip_ref.namelist()
                for file_name in file_list:
                    # Skip macOS metadata files and empty directory entries
                    if file_name.startswith('__MACOSX') or file_name.endswith('/'):
                        continue

                    logger.info(f"Extracting and processing: {file_name}")
                    with zip_ref.open(file_name) as extracted_file:
                        extracted_content = extracted_file.read()
                        file_extension = os.path.splitext(file_name)[1].lower()
                        
                        # Create a temporary URL-like identifier for caching
                        # This helps reuse the parsing logic and potential caching
                        temp_url_for_parsing = f"zip://{file_name}"

                        # Recursively process the extracted file
                        # We create a new task to process each file to avoid blocking
                        # This demonstrates a more advanced, but powerful, pattern
                        text, metadata = await self.get_text_and_metadata(
                            extracted_content, temp_url_for_parsing
                        )

                        if text.strip():
                            combined_text.append(f"--- START OF FILE: {file_name} ---\n{text}\n--- END OF FILE: {file_name} ---")
                            combined_metadata.extend(metadata)
        
        return "\n\n".join(combined_text), combined_metadata    

    # ... (inside the HybridFastTrackRAGPipeline class)

    def _parse_image(self, content: bytes) -> Tuple[str, List[Dict]]:
        """Parse image files using OCR"""
        if not IMAGE_PROCESSING_AVAILABLE:
            logger.warning("Image processing libraries not installed. Skipping OCR.")
            return "Image content could not be processed.", [{'type': 'image_error'}]
        try:
            image = Image.open(io.BytesIO(content))
            text = pytesseract.image_to_string(image)
            return self._clean_text(text), [{'type': 'image'}]
        except Exception as e:
            logger.error(f"OCR processing failed: {e}")
            return "Unable to extract text from image.", [{'type': 'ocr_error'}]

# ... (rest of the class)    
        
    async def download_document_async(self, url: str) -> bytes:
        """Async document download with better error handling"""
        headers = {
            'User-Agent': 'Mozilla/5.0 (compatible; RAGPipeline/3.0)',
            'Accept': '*/*'
        }
        
        try:
            async with aiohttp.ClientSession() as session:
                async with session.get(url, headers=headers, 
                                     timeout=aiohttp.ClientTimeout(total=30)) as response:
                    response.raise_for_status()
                    content = await response.read()
                    logger.info(f"Downloaded {len(content)/1024/1024:.1f}MB from {url}")
                    return content
        except Exception as e:
            logger.error(f"Async download failed: {e}, trying sync download")
            # Fallback to sync download
            return await asyncio.get_event_loop().run_in_executor(
                None, self._download_sync, url
            )
    
    def _download_sync(self, url: str) -> bytes:
        """Sync download fallback"""
        try:
            response = requests.get(url, timeout=30, stream=True)
            response.raise_for_status()
            
            content = b""
            for chunk in response.iter_content(chunk_size=1024*1024):
                content += chunk
                if len(content) > 200*1024*1024:  # 200MB limit
                    break
            
            return content
        except Exception as e:
            logger.error(f"Sync download also failed: {e}")
            raise
    
    def _parse_pdf_enhanced(self, content: bytes) -> Tuple[str, List[Dict]]:
        """Enhanced PDF parsing with multiple fallbacks"""
        logger.info("Starting enhanced PDF parsing")
        full_text = ""
        chunk_metadata = []
        
        # Try pdfplumber first
        try:
            temp_file = io.BytesIO(content)
            with pdfplumber.open(temp_file) as pdf:
                total_pages = len(pdf.pages)
                logger.info(f"Processing {total_pages} pages with pdfplumber")
                
                for i, page in enumerate(pdf.pages):
                    if i % 50 == 0 and i > 0:
                        logger.info(f"Processing page {i}/{total_pages}")
                    
                    # Extract text
                    page_text = page.extract_text() or ""
                    if page_text.strip():
                        full_text += f"\n--- PAGE {i+1} ---\n{page_text}\n"
                        chunk_metadata.append({
                            'page': i+1,
                            'type': 'text'
                        })
                    
                    # Extract tables (limit to avoid memory issues)
                    if i < 100:  # Only extract tables from first 100 pages
                        tables = page.extract_tables()
                        for j, table in enumerate(tables[:3]):  # Max 3 tables per page
                            if table:
                                table_text = self._format_table(table)
                                if table_text:
                                    full_text += f"\n=== TABLE {j+1} (Page {i+1}) ===\n{table_text}\n"
                                    chunk_metadata.append({
                                        'page': i+1,
                                        'type': 'table',
                                        'table_idx': j+1
                                    })
                
                logger.info(f"Successfully parsed {total_pages} pages")
                return self._clean_text(full_text), chunk_metadata
               
        except Exception as e:
            logger.warning(f"pdfplumber failed: {e}, trying pypdf")
        
        # Fallback to pypdf
        try:
            temp_file = io.BytesIO(content)
            reader = pypdf.PdfReader(temp_file)
            total_pages = len(reader.pages)
            logger.info(f"Processing {total_pages} pages with pypdf")
            
            for i, page in enumerate(reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text:
                        full_text += f"\n--- PAGE {i+1} ---\n{page_text}\n"
                        chunk_metadata.append({
                            'page': i+1,
                            'type': 'text_fallback'
                        })
                except Exception as e:
                    logger.warning(f"Page {i+1} extraction failed: {e}")
            
            return self._clean_text(full_text), chunk_metadata
            
        except Exception as e:
            logger.error(f"All PDF parsing methods failed: {e}")
            raise ValueError("Could not parse PDF document")
    
    def _format_table(self, table: List[List]) -> str:
        """Format table with better structure"""
        if not table:
            return ""
        
        formatted_rows = []
        
        # Add header separator if first row looks like headers
        if table and len(table) > 1:
            first_row = table[0]
            if all(cell and isinstance(cell, str) for cell in first_row):
                formatted_rows.append(" | ".join(str(cell) for cell in first_row))
                formatted_rows.append("-" * min(len(formatted_rows[0]), 100))
                
                for row in table[1:]:
                    if any(cell for cell in row):
                        formatted_rows.append(" | ".join(str(cell) if cell else "" for cell in row))
            else:
                for row in table:
                    if any(cell for cell in row):
                        formatted_rows.append(" | ".join(str(cell) if cell else "" for cell in row))
        
        return "\n".join(formatted_rows)
    
    def _clean_text(self, text: str) -> str:
        """Enhanced text cleaning"""
        # Remove excessive whitespace
        text = re.sub(r'\n\s*\n\s*\n', '\n\n', text)
        text = re.sub(r'[ \t]+', ' ', text)
        
        # Fix common OCR issues
        replacements = {
            'ﬁ': 'fi', 'ﬂ': 'fl', '–': '-', '"': '"', '"': '"',
            ''': "'", ''': "'", '…': '...', '—': '--'
        }
        for old, new in replacements.items():
            text = text.replace(old, new)
        
        # Fix broken words
        text = re.sub(r'(\w+)-\s*\n\s*(\w+)', r'\1\2', text)
        
        return text.strip()
    
    def _parse_docx(self, temp_file: io.BytesIO) -> str:
        """Enhanced DOCX parsing with better structure preservation"""
        try:
            doc = Document(temp_file)
            full_text = []
            
            # Extract document properties if available
            if doc.core_properties.title:
                full_text.append(f"Title: {doc.core_properties.title}")
            if doc.core_properties.subject:
                full_text.append(f"Subject: {doc.core_properties.subject}")
            
            full_text.append("\n=== DOCUMENT CONTENT ===\n")
            
            # Process paragraphs with style information
            for para in doc.paragraphs:
                if para.text.strip():
                    # Detect headings
                    if para.style.name.startswith('Heading'):
                        level = para.style.name[-1] if para.style.name[-1].isdigit() else '1'
                        full_text.append(f"\n{'#' * int(level)} {para.text}\n")
                    else:
                        full_text.append(para.text)
            
            # Extract tables with better formatting
            for i, table in enumerate(doc.tables):
                full_text.append(f"\n=== TABLE {i+1} ===")
                
                # Get headers if first row looks like headers
                headers = []
                if table.rows:
                    first_row = table.rows[0]
                    headers = [cell.text.strip() for cell in first_row.cells]
                    if all(h for h in headers):
                        full_text.append("Headers: " + " | ".join(headers))
                
                # Process remaining rows
                for j, row in enumerate(table.rows[1:] if headers else table.rows):
                    row_data = []
                    for k, cell in enumerate(row.cells):
                        cell_text = cell.text.strip()
                        if headers and k < len(headers):
                            row_data.append(f"{headers[k]}: {cell_text}")
                        else:
                            row_data.append(cell_text)
                    
                    if any(row_data):
                        full_text.append(f"Row {j+1}: " + " | ".join(row_data))
                
                full_text.append("=== END TABLE ===\n")
            
            return "\n".join(full_text)
            
        except Exception as e:
            logger.error(f"DOCX parsing error: {e}")
            return f"Error parsing DOCX: {str(e)}"
    
    def _parse_odt(self, temp_file: io.BytesIO) -> str:
        """Parse ODT files"""
        try:
            doc = load(temp_file)
            full_text = []
            
            for element in doc.getElementsByType(P):
                text = str(element)
                if text.strip():
                    full_text.append(text)
            
            return "\n\n".join(full_text)
        except Exception as e:
            logger.error(f"ODT parsing error: {e}")
            return ""
    
    def _parse_other_formats(self, content: bytes, file_extension: str) -> Tuple[str, List[Dict]]:
        """Parse non-PDF formats with better handling"""
        logger.info(f"Parsing document with extension: {file_extension}")
        
        # Skip binary formats that can't contain text
        binary_extensions = ['.gif', '.bmp', '.zip', '.rar', 
                            '.7z', '.mp3', '.mp4', '.avi', '.mov', '.bin']
        if file_extension.lower() in binary_extensions:
            logger.warning(f"Skipping binary file format: {file_extension}")
            return "This appears to be a binary file that cannot be processed for text content.", [{'type': 'binary_skip'}]
        
        # Handle Excel files
        # In app/core/rag_pipeline.py - Enhanced Excel parsing
        if file_extension in ['.xlsx', '.xls', '.csv']:
            try:
                temp_file = io.BytesIO(content)
                
                # For CSV files
                if file_extension == '.csv':
                    import csv
                    text_content = content.decode('utf-8', errors='ignore')
                    reader = csv.DictReader(io.StringIO(text_content))
                    
                    # Create structured text with headers
                    headers = reader.fieldnames
                    text = f"Column Headers: {', '.join(headers)}\n\n"
                    
                    # Add data rows with context
                    rows_data = []
                    for i, row in enumerate(reader):
                        if i < 1000:  # Limit for performance
                            row_text = " | ".join([f"{k}: {v}" for k, v in row.items() if v])
                            rows_data.append(f"Row {i+1}: {row_text}")
                    
                    text += "\n".join(rows_data)
                    
                    # Also create a summary section
                    df = pd.read_csv(io.BytesIO(content))
                    text += f"\n\n=== DATA SUMMARY ===\n"
                    text += f"Total Rows: {len(df)}\n"
                    text += f"Columns: {', '.join(df.columns)}\n"
                    for col in df.columns:
                        if df[col].dtype in ['int64', 'float64']:
                            text += f"{col} - Min: {df[col].min()}, Max: {df[col].max()}, Mean: {df[col].mean():.2f}\n"
                        else:
                            text += f"{col} - Unique values: {df[col].nunique()}\n"
                            if df[col].nunique() < 20:
                                text += f"  Values: {', '.join(df[col].unique()[:20].astype(str))}\n"
                
                else:  # Excel files
                    df = pd.read_excel(temp_file, engine='openpyxl' if file_extension == '.xlsx' else None)
                    
                    # Convert with better structure
                    text = f"=== EXCEL DATA ===\n"
                    text += f"Shape: {df.shape[0]} rows × {df.shape[1]} columns\n"
                    text += f"Columns: {', '.join(df.columns)}\n\n"
                    
                    # Add full data in searchable format
                    for idx, row in df.iterrows():
                        if idx < 1000:  # Limit
                            row_text = " | ".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val)])
                            text += f"Row {idx+1}: {row_text}\n"
                    
                    # Add statistical summary for numeric columns
                    text += "\n=== STATISTICS ===\n"
                    text += df.describe().to_string()
                    
                return self._clean_text(text), [{'type': 'spreadsheet', 'headers': list(df.columns)}]
                
            except Exception as e:
                logger.error(f"Failed to parse Excel/CSV file: {e}")
                return f"Error parsing file: {str(e)}", [{'type': 'excel_error'}]
        
        # Handle PowerPoint
        # In _parse_other_formats, enhance PowerPoint section:
        if file_extension in ['.pptx', '.ppt'] and PPTX_AVAILABLE:
            try:
                temp_file = io.BytesIO(content)
                prs = Presentation(temp_file)
                text_runs = []
                
                # Add presentation title if available
                if prs.core_properties.title:
                    text_runs.append(f"Presentation Title: {prs.core_properties.title}\n")
                
                for slide_num, slide in enumerate(prs.slides):
                    slide_text = f"\n=== SLIDE {slide_num + 1} ===\n"
                    
                    # Extract title
                    if slide.shapes.title:
                        slide_text += f"Title: {slide.shapes.title.text}\n"
                    
                    # Extract all text from shapes
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            # Check if it's a table
                            if shape.has_table:
                                slide_text += "Table:\n"
                                for row in shape.table.rows:
                                    row_text = " | ".join([cell.text for cell in row.cells])
                                    slide_text += f"  {row_text}\n"
                            else:
                                slide_text += f"{shape.text}\n"
                        
                        # Extract text from text frame
                        if hasattr(shape, "text_frame"):
                            for paragraph in shape.text_frame.paragraphs:
                                if paragraph.text.strip():
                                    slide_text += f"- {paragraph.text}\n"
                    
                    text_runs.append(slide_text)
                
                text = "\n".join(text_runs)
                return self._clean_text(text), [{'type': 'presentation'}]
                
            except Exception as e:
                logger.warning(f"Failed to parse PowerPoint: {e}")
                return f"PowerPoint parsing error: {str(e)}", [{'type': 'pptx_error'}]
        # Handle DOCX
        if file_extension in ['.docx', '.doc']:
            try:
                temp_file = io.BytesIO(content)
                text = self._parse_docx(temp_file)
                return self._clean_text(text), [{'type': 'docx'}]
            except Exception as e:
                logger.warning(f"Failed to parse DOCX: {e}")
        
        # Handle ODT
        if file_extension == '.odt':
            try:
                temp_file = io.BytesIO(content)
                text = self._parse_odt(temp_file)
                return self._clean_text(text), [{'type': 'odt'}]
            except Exception as e:
                logger.warning(f"Failed to parse ODT: {e}")
        
        # Default to text extraction
        try:
            text = content.decode('utf-8', errors='ignore')
            return self._clean_text(text), [{'type': 'text'}]
        except Exception as e:
            logger.error(f"Failed to decode as text: {e}")
            return "Unable to extract text from this file format.", [{'type': 'unknown'}]
    
    async def get_or_create_optimized_vector_store(self, url: str) -> OptimizedVectorStore:
        """Create vector store with enhanced caching"""
        # Generate cache key
        cache_key = f"vs_{hashlib.md5(url.encode()).hexdigest()}"
        
        # Check cache
        cached_store = await cache.get(cache_key)
        if cached_store:
            logger.info(f"Using cached vector store for {url}")
            return cached_store
        
        # Check for cached intermediate embeddings
        embed_cache_key = f"emb_{hashlib.md5(url.encode()).hexdigest()}"
        cached_embeddings = await cache.get(embed_cache_key)
        
        if cached_embeddings:
            logger.info("Using cached embeddings")
            chunks, embeddings, chunk_metadata = cached_embeddings
        else:
            # Download and parse document
            logger.info(f"Downloading document from {url}")
            try:
                content = await self.download_document_async(url)
            except Exception as e:
                logger.error(f"Download failed: {e}")
                raise ValueError(f"Could not download document: {str(e)}")
            
            # Determine file type and parse
            file_extension = os.path.splitext(url.split('?')[0])[1].lower()

            if file_extension == '.zip':
                text, doc_metadata = await self._parse_zip(content)
            else:
                # Use the new helper for all other file types
                text, doc_metadata = await self.get_text_and_metadata(content, url)
            
            
            
            # Check if we got valid text
            if not text or len(text) < 10:
                logger.warning(f"No valid text extracted from {url}")
                text = "Document appears to be empty or could not be parsed."
                doc_metadata = [{'type': 'empty'}]
            
            # Smart chunking
            chunks, chunk_metadata = SmartChunker.chunk_with_boundaries(text, doc_metadata)
            
            if not chunks:
                chunks = [text]  # Fallback to full text as single chunk
                chunk_metadata = doc_metadata
            
            # Add context to chunks for better retrieval
            contextualized_chunks = ContextualRetrieval.add_context_to_chunks(chunks, chunk_metadata)
            
            logger.info(f"Created {len(chunks)} chunks, generating embeddings...")
            
            # Generate embeddings in batches
            batch_size = 32
            all_embeddings = []
            
            for i in range(0, len(contextualized_chunks), batch_size):
                batch = contextualized_chunks[i:i + batch_size]
                try:
                    batch_embeddings = await asyncio.wait_for(
                        asyncio.get_event_loop().run_in_executor(
                            None,
                            lambda b: self.embedding_model.encode(b, convert_to_numpy=True, show_progress_bar=False),
                            batch
                        ),
                        timeout=30.0
                    )
                    all_embeddings.append(batch_embeddings)
                except asyncio.TimeoutError:
                    logger.error(f"Embedding generation timeout for batch {i//batch_size}")
                    # Use zero embeddings as fallback
                    zero_embeddings = np.zeros((len(batch), self.embedding_model.get_sentence_embedding_dimension()))
                    all_embeddings.append(zero_embeddings)
            
            embeddings = np.vstack(all_embeddings).astype('float32')
            
            # Cache embeddings
            await cache.set(embed_cache_key, (chunks, embeddings, chunk_metadata))
            logger.info(f"Cached embeddings for {url}")
        
        # Create vector store
        vector_store = OptimizedVectorStore(chunks, embeddings, self.embedding_model, chunk_metadata)
        
        # Cache the store
        await cache.set(cache_key, vector_store)
        
        return vector_store
    
    async def _answer_question_fast(self, question: str, vector_store: OptimizedVectorStore) -> str:
        """Enhanced answer generation with better context understanding"""
        
        # Check cache first
        answer_cache_key = f"ans_{hashlib.md5(question.encode()).hexdigest()[:16]}"
        cached_answer = await cache.get(answer_cache_key)
        if cached_answer:
            return cached_answer
        
        # Detect if question needs aggregation or calculation
        needs_aggregation = any(word in question.lower() for word in [
            'how many', 'count', 'total', 'sum', 'average', 'highest', 'lowest',
            'maximum', 'minimum', 'list all', 'compare', 'ratio', 'percentage'
        ])
        
        # Get more chunks for aggregation questions
        k = 20 if needs_aggregation else settings.MAX_CHUNKS_PER_QUERY
        retrieved_results = vector_store.hybrid_search(question, k=k)
        
        if not retrieved_results:
            return "Information not found in the provided context."
        
        # Use more chunks for complex questions
        num_chunks = 12 if needs_aggregation else 8
        top_chunks = [result[0] for result in retrieved_results[:num_chunks]]
        
        # Enhanced prompt based on question type
        if needs_aggregation:
            prompt = f"""You are analyzing data from a document. Answer the question by carefully examining ALL the provided context.

    CONTEXT:
    {chr(10).join(top_chunks)}

    QUESTION: {question}

    INSTRUCTIONS:
    - For counting/listing questions: Examine ALL occurrences in the context
    - For comparison questions: Compare the specific items mentioned
    - For statistical questions: Calculate based on the data provided
    - Include specific numbers, names, or values found
    - If data is incomplete, mention what you found and what might be missing
    - Be precise and comprehensive

    ANSWER:"""
        else:
            prompt = f"""Answer the question based on the provided context. Be specific and direct.

    CONTEXT:
    {chr(10).join(top_chunks)}

    QUESTION: {question}

    INSTRUCTIONS:
    - Extract the specific information requested
    - Include exact values, terms, or details from the context
    - If partially available, provide what you found
    - Say "Information not found" only if truly not present

    ANSWER:"""
        
        try:
            # Use more capable model for complex questions
            model_name = settings.LLM_MODEL_NAME_PRECISE if needs_aggregation else settings.LLM_MODEL_NAME
            model = genai.GenerativeModel(model_name)
            
            response = await asyncio.wait_for(
                model.generate_content_async(
                    prompt,
                    generation_config=genai.types.GenerationConfig(
                        temperature=0.1,
                        max_output_tokens=600 if needs_aggregation else 400,
                        top_p=0.95,
                        candidate_count=1
                    )
                ),
                timeout=settings.ANSWER_TIMEOUT_SECONDS * (1.5 if needs_aggregation else 1)
            )
            
            answer = response.text.strip()
            
            # Better validation
            if not answer or len(answer) < 5:
                answer = "Unable to generate answer from the document."
            elif "information not found" in answer.lower() and len(top_chunks) > 0:
                # Try once more with a simpler prompt
                simple_prompt = f"Based on this text, {question}\n\nText: {' '.join(top_chunks[:3])}"
                try:
                    response = await model.generate_content_async(simple_prompt)
                    fallback_answer = response.text.strip()
                    if fallback_answer and len(fallback_answer) > 10:
                        answer = fallback_answer
                except:
                    pass
            
            # Cache successful answers
            if len(answer) > 20 and "error" not in answer.lower():
                await cache.set(answer_cache_key, answer)
            
            return answer
            
        except Exception as e:
            logger.error(f"Answer generation failed: {e}")
            return "Unable to process this question."
    
    async def process_query(self, document_url: str, questions: List[str]) -> List[str]:
        """Process queries with optimized parallel execution"""
        start_time = time.time()
        logger.info(f"Processing {len(questions)} questions for document: {document_url}")
        
        try:
            # Create vector store once
            vector_store = await self.get_or_create_optimized_vector_store(document_url)
            
            # Process all questions in parallel with controlled concurrency
            semaphore = asyncio.Semaphore(settings.MAX_CONCURRENT_QUESTIONS)
            
            async def process_with_semaphore(q):
                async with semaphore:
                    return await self._answer_question_safe(q, vector_store)
            
            # Create all tasks
            tasks = [process_with_semaphore(q) for q in questions]
            
            # Execute with overall timeout
            try:
                answers = await asyncio.wait_for(
                    asyncio.gather(*tasks, return_exceptions=True),
                    timeout=settings.TOTAL_TIMEOUT_SECONDS
                )
                
                # Process results
                final_answers = []
                for i, answer in enumerate(answers):
                    if isinstance(answer, Exception):
                        logger.error(f"Question {i+1} error: {answer}")
                        final_answers.append("Error processing this question.")
                    else:
                        final_answers.append(answer)
                
                processing_time = time.time() - start_time
                logger.info(f"Completed {len(questions)} questions in {processing_time:.2f}s "
                            f"(avg: {processing_time/len(questions):.2f}s per question)")
                
                return final_answers
                
            except asyncio.TimeoutError:
                logger.error("Overall processing timeout")
                return ["Processing timeout. Please try again."] * len(questions)
                
        except Exception as e:
            logger.critical(f"Critical error: {e}", exc_info=True)
            return ["Document processing error."] * len(questions)
    
    async def _answer_question_safe(self, question: str, vector_store: OptimizedVectorStore) -> str:
        """Safe wrapper for question answering"""
        try:
            return await self._answer_question_fast(question, vector_store)
        except Exception as e:
            logger.error(f"Error answering question: {e}")
            return "Unable to process this question due to an error."
    
    async def process_query_with_explainability(self, document_url: str, questions: List[str]) -> Tuple[List[str], List[Dict], float]:
        """Process with detailed explanations (backward compatibility)"""
        start_time = time.time()
        
        # Use the fast processing method
        simple_answers = await self.process_query(document_url, questions)
        
        # Create basic detailed answers
        detailed_answers = []
        for answer in simple_answers:
            detailed_answer = {
                "answer": answer,
                "confidence": 0.8 if answer and "error" not in answer.lower() else 0.0,
                "source_clauses": [],
                "reasoning": "Answer extracted using optimized hybrid retrieval.",
                "coverage_decision": self._determine_coverage(answer)
            }
            detailed_answers.append(detailed_answer)
        
        processing_time = time.time() - start_time
        return simple_answers, detailed_answers, processing_time
    
    def _determine_coverage(self, answer: str) -> str:
        """Determine coverage decision"""
        answer_lower = answer.lower()
        
        if any(term in answer_lower for term in ['not covered', 'excluded', 'not available', 'cannot find', 'not found']):
            return "Not Covered"
        elif any(term in answer_lower for term in ['covered', 'included', 'eligible', 'yes']):
            return "Covered"
        elif any(term in answer_lower for term in ['conditions apply', 'subject to', 'depending on']):
            return "Conditional"
        elif any(term in answer_lower for term in ['error', 'timeout', 'unable']):
            return "Error"
        else:
            return "Review Required"