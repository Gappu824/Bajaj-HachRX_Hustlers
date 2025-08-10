# app/core/document_parser.py - Comprehensive document parsing with fixes
import io
import re
import logging
import zipfile
import tempfile
import shutil
import os
from typing import Tuple, List, Dict, Optional
import pandas as pd
import numpy as np
from PIL import Image
import pytesseract
import pdfplumber
import pypdf
from docx import Document
from pptx import Presentation
from odf.text import P
from odf.opendocument import load
from app.core.smart_chunker import SmartChunker
from app.core.config import settings
import asyncio

logger = logging.getLogger(__name__)

class DocumentParser:
    """Unified document parser with format-specific optimizations"""

    @staticmethod
    async def parse_bin_incrementally(file_path: str, vector_store, pipeline):
        """Stream-process a large binary file"""
        READ_CHUNK_SIZE = 1024 * 1024  # 1MB chunks
        
        try:
            with open(file_path, 'rb') as f:
                while True:
                    byte_chunk = f.read(READ_CHUNK_SIZE)
                    if not byte_chunk:
                        break
                    
                    text, _ = DocumentParser.parse_binary(byte_chunk)
                    if not text.strip():
                        continue
                    
                    chunks, chunk_meta = SmartChunker.chunk_document(
                        text, 
                        [{'type': 'binary_stream'}],
                        chunk_size=settings.CHUNK_SIZE_CHARS,
                        overlap=settings.CHUNK_OVERLAP_CHARS
                    )
                    
                    if chunks:
                        embeddings = await pipeline._generate_embeddings(chunks)
                        vector_store.add(chunks, embeddings, chunk_meta)
                        
        except Exception as e:
            logger.error(f"Failed to process binary file: {e}")
    
    @staticmethod
    def parse_document(content: bytes, file_extension: str) -> Tuple[str, List[Dict]]:
        """Route to appropriate parser with fallback to universal parser"""
        
        # Normalize extension
        file_extension = file_extension.lower()
        if not file_extension.startswith('.'):
            file_extension = '.' + file_extension
        
        # Define the primary parser mapping
        parser_map = {
            '.pdf': DocumentParser.parse_pdf,
            '.xlsx': lambda c: DocumentParser.parse_spreadsheet(c, '.xlsx'),
            '.xls': lambda c: DocumentParser.parse_spreadsheet(c, '.xls'),
            '.csv': lambda c: DocumentParser.parse_spreadsheet(c, '.csv'),
            '.docx': DocumentParser.parse_word,
            '.doc': DocumentParser.parse_word,
            '.pptx': DocumentParser.parse_powerpoint,
            '.ppt': DocumentParser.parse_powerpoint,
            '.png': DocumentParser.parse_image,
            '.jpg': DocumentParser.parse_image,
            '.jpeg': DocumentParser.parse_image,
            '.tiff': DocumentParser.parse_image,
            '.bmp': DocumentParser.parse_image,
            '.odt': DocumentParser.parse_odt,
            '.bin': DocumentParser.parse_binary,
        }
        
        # Get primary parser or default to text parser
        primary_parser = parser_map.get(file_extension, DocumentParser.parse_text)
        
        try:
            # Attempt primary parsing
            text, metadata = primary_parser(content)

            # Validate parsing results
            if (not text or 
                len(text.strip()) < 20 or 
                "unable to parse" in text.lower() or 
                "error parsing" in text.lower()):
                
                logger.warning(f"Primary parser for '{file_extension}' returned insufficient content")
                raise ValueError("Primary parser yielded insufficient content")

            return text, metadata

        except Exception as e:
            logger.info(f"Primary parser failed: {e}. Attempting universal parser fallback")
            
            try:
                # Import here to avoid circular dependency
                from app.core.universal_parser import UniversalDocumentParser
                
                universal_parser = UniversalDocumentParser()
                return universal_parser.parse_any_document(content, f"file{file_extension}")
                
            except Exception as universal_e:
                logger.error(f"All parsing methods failed for {file_extension}: {universal_e}")
                return (
                    f"Unable to process document. Format '{file_extension}' parsing failed: {str(universal_e)[:200]}", 
                    [{'type': 'fatal_error', 'extension': file_extension}]
                )
    
    @staticmethod
    def parse_binary(content: bytes) -> Tuple[str, List[Dict]]:
        """Extract printable strings from binary content"""
        try:
            # Attempt UTF-8 decoding with replacement
            text = content.decode('utf-8', errors='replace')
            
            # Extract sequences of printable characters (4+ chars)
            printable_strings = re.findall(r'[a-zA-Z0-9\s\.\-\_]{4,}', text)
            
            if printable_strings:
                extracted_text = " ".join(printable_strings[:500])  # Limit output
                logger.info(f"Extracted {len(extracted_text)} characters from binary file")
                return extracted_text, [{'type': 'binary_extract', 'strings_found': len(printable_strings)}]
            else:
                return "No meaningful text found in binary file", [{'type': 'binary_extract'}]

        except Exception as e:
            logger.error(f"Binary parsing failed: {e}")
            return "Unable to process binary file", [{'type': 'error'}]
    
    @staticmethod
    def parse_pdf(content: bytes) -> Tuple[str, List[Dict]]:
        """Enhanced PDF parsing with table extraction"""
        full_text = []
        metadata = []
        
        try:
            # Primary: pdfplumber for best text and table extraction
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                for i, page in enumerate(pdf.pages):
                    # Extract text
                    page_text = page.extract_text() or ""
                    if page_text.strip():
                        full_text.append(f"\n=== PAGE {i+1} ===\n{page_text}")
                        metadata.append({'page': i+1, 'type': 'text'})
                    
                    # Extract tables
                    tables = page.extract_tables()
                    for j, table in enumerate(tables):
                        if table and len(table) > 1:
                            table_text = DocumentParser._format_table(table)
                            full_text.append(f"\n=== TABLE {j+1} (Page {i+1}) ===\n{table_text}")
                            metadata.append({'page': i+1, 'type': 'table', 'table_id': j+1})
        
        except Exception as e:
            logger.warning(f"pdfplumber failed: {e}, trying pypdf fallback")
            # Fallback: pypdf
            try:
                reader = pypdf.PdfReader(io.BytesIO(content))
                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        full_text.append(f"\n=== PAGE {i+1} ===\n{page_text}")
                        metadata.append({'page': i+1, 'type': 'text_fallback'})
            except Exception as e2:
                logger.error(f"All PDF parsing methods failed: {e2}")
                return "Unable to parse PDF document", [{'type': 'error'}]
        
        final_text = "\n".join(full_text) if full_text else "No text extracted from PDF"
        return final_text, metadata or [{'type': 'empty_pdf'}]
    
    @staticmethod
    def parse_spreadsheet(content: bytes, file_extension: str) -> Tuple[str, List[Dict]]:
        """Enhanced spreadsheet parsing with structure preservation"""
        try:
            # Read the spreadsheet
            if file_extension == '.csv':
                df = pd.read_csv(io.BytesIO(content), encoding='utf-8', errors='ignore')
            else:
                engine = 'openpyxl' if file_extension == '.xlsx' else 'xlrd'
                df = pd.read_excel(io.BytesIO(content), engine=engine)
            
            text_parts = []
            
            # 1. Summary information
            text_parts.append("=== SPREADSHEET SUMMARY ===")
            text_parts.append(f"Rows: {len(df)}")
            text_parts.append(f"Columns: {len(df.columns)}")
            text_parts.append(f"Column Names: {', '.join(str(col) for col in df.columns)}")
            
            # 2. Numeric column statistics
            numeric_cols = df.select_dtypes(include=[np.number]).columns
            if len(numeric_cols) > 0:
                text_parts.append("\n=== NUMERIC STATISTICS ===")
                for col in numeric_cols[:10]:  # Limit to 10 columns
                    try:
                        stats = df[col].describe()
                        text_parts.append(f"{col}: Min={stats['min']:.2f}, Max={stats['max']:.2f}, Mean={stats['mean']:.2f}")
                    except:
                        pass
            
            # 3. Categorical information
            cat_cols = df.select_dtypes(include=['object']).columns
            if len(cat_cols) > 0:
                text_parts.append("\n=== CATEGORICAL DATA ===")
                for col in cat_cols[:10]:  # Limit to 10 columns
                    unique_count = df[col].nunique()
                    text_parts.append(f"{col}: {unique_count} unique values")
                    if unique_count <= 20:
                        values = df[col].dropna().unique()[:20]
                        text_parts.append(f"  Values: {', '.join(str(v) for v in values)}")
            
            # 4. Data sample
            text_parts.append("\n=== DATA ROWS ===")
            text_parts.append(" | ".join(str(col) for col in df.columns))
            
            # Smart sampling for different dataset sizes
            if len(df) <= 100:
                # Small dataset: include all rows
                for idx, row in df.iterrows():
                    row_text = " | ".join(str(val) if pd.notna(val) else "" for val in row.values)
                    text_parts.append(f"Row {idx+1}: {row_text}")
            elif len(df) <= 1000:
                # Medium dataset: first 50, last 50
                for idx, row in df.head(50).iterrows():
                    row_text = " | ".join(str(val) if pd.notna(val) else "" for val in row.values)
                    text_parts.append(f"Row {idx+1}: {row_text}")
                
                text_parts.append(f"\n... ({len(df) - 100} rows omitted) ...\n")
                
                for idx, row in df.tail(50).iterrows():
                    row_text = " | ".join(str(val) if pd.notna(val) else "" for val in row.values)
                    text_parts.append(f"Row {idx+1}: {row_text}")
            else:
                # Large dataset: sample strategically
                sample_size = min(200, len(df) // 10)
                sample_df = df.sample(n=sample_size)
                for idx, row in sample_df.iterrows():
                    row_text = " | ".join(str(val) if pd.notna(val) else "" for val in row.values)
                    text_parts.append(f"Row {idx+1}: {row_text}")
            
            metadata = [{
                'type': 'spreadsheet', 
                'rows': len(df), 
                'columns': len(df.columns),
                'file_extension': file_extension
            }]
            
            return "\n".join(text_parts), metadata
            
        except Exception as e:
            logger.error(f"Spreadsheet parsing failed: {e}")
            return f"Error parsing spreadsheet: {str(e)}", [{'type': 'error'}]
    
    @staticmethod
    def parse_word(content: bytes) -> Tuple[str, List[Dict]]:
        """Enhanced Word document parsing with pypdf fallback"""
        try:
            doc = Document(io.BytesIO(content))
            text_parts = []
            metadata = []
            
            # Document properties
            try:
                if doc.core_properties.title:
                    text_parts.append(f"Title: {doc.core_properties.title}")
                if doc.core_properties.subject:
                    text_parts.append(f"Subject: {doc.core_properties.subject}")
            except:
                pass
            
            # Process paragraphs with style preservation
            for para in doc.paragraphs:
                if para.text and para.text.strip():
                    try:
                        if para.style.name.startswith('Heading'):
                            level = para.style.name[-1] if para.style.name[-1].isdigit() else '1'
                            text_parts.append(f"\n{'#' * int(level)} {para.text}")
                            metadata.append({'type': 'heading', 'level': int(level)})
                        else:
                            text_parts.append(para.text)
                            metadata.append({'type': 'paragraph'})
                    except:
                        text_parts.append(para.text)
                        metadata.append({'type': 'paragraph'})
            
            # Process tables
            for i, table in enumerate(doc.tables):
                text_parts.append(f"\n=== TABLE {i+1} ===")
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    if any(row_text):
                        text_parts.append(" | ".join(row_text))
                text_parts.append("=== END TABLE ===\n")
                metadata.append({'type': 'table', 'table_id': i+1})

            if text_parts:
                return "\n".join(text_parts), metadata
            else:
                raise ValueError("No text extracted from Word document")
                
        except Exception as e:
            logger.warning(f"python-docx failed: {e}. Trying pypdf fallback")
            try:
                reader = pypdf.PdfReader(io.BytesIO(content))
                text_parts = []
                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        text_parts.append(page_text)
                
                if text_parts:
                    return "\n\n".join(text_parts), [{'type': 'text_fallback'}]
                else:
                    return "Unable to extract text from Word document", [{'type': 'error'}]
                    
            except Exception as e2:
                logger.error(f"All Word parsing methods failed: {e2}")
                return "Error parsing Word document", [{'type': 'error'}]
    
    @staticmethod
    def parse_powerpoint(content: bytes) -> Tuple[str, List[Dict]]:
        """Enhanced PowerPoint parsing"""
        try:
            prs = Presentation(io.BytesIO(content))
            text_parts = []
            metadata = []
            
            # Presentation title
            try:
                if prs.core_properties.title:
                    text_parts.append(f"Presentation: {prs.core_properties.title}\n")
            except:
                pass
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"\n=== SLIDE {slide_num} ===")
                slide_content = []
                
                # Title
                try:
                    if slide.shapes.title and slide.shapes.title.text:
                        slide_content.append(f"Title: {slide.shapes.title.text}")
                except:
                    pass
                
                # Process all shapes
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, "text") and shape.text and shape != slide.shapes.title:
                            slide_content.append(shape.text)
                        
                        # Tables
                        if hasattr(shape, 'has_table') and shape.has_table:
                            slide_content.append("Table:")
                            for row in shape.table.rows:
                                row_text = " | ".join([cell.text for cell in row.cells])
                                slide_content.append(f"  {row_text}")
                    except:
                        continue
                
                text_parts.extend(slide_content)
                metadata.append({'type': 'slide', 'slide_num': slide_num})
            
            return "\n".join(text_parts), metadata
            
        except Exception as e:
            logger.error(f"PowerPoint parsing failed: {e}")
            return f"Error parsing PowerPoint: {str(e)}", [{'type': 'error'}]
    
    @staticmethod
    def parse_image(content: bytes) -> Tuple[str, List[Dict]]:
        """Parse image using OCR with error handling"""
        try:
            image = Image.open(io.BytesIO(content))
            
            # OCR text extraction
            text = pytesseract.image_to_string(image)
            
            # Image metadata
            metadata = [{
                'type': 'image',
                'format': getattr(image, 'format', 'unknown'),
                'size': getattr(image, 'size', (0, 0)),
                'mode': getattr(image, 'mode', 'unknown')
            }]
            
            return text.strip() if text.strip() else "No text found in image", metadata
                
        except Exception as e:
            logger.error(f"Image parsing failed: {e}")
            return "Unable to process image", [{'type': 'error'}]
    
    @staticmethod
    async def parse_zip_incrementally(zip_path: str, vector_store, pipeline):
        """Extract ZIP and process contents incrementally"""
        temp_dir = tempfile.mkdtemp()
        
        try:
            with zipfile.ZipFile(zip_path, 'r') as zf:
                zf.extractall(temp_dir)

            # Process files incrementally
            processed_count = 0
            max_files = 50  # Limit to prevent overwhelming
            
            for root, _, files in os.walk(temp_dir):
                for filename in files:
                    if processed_count >= max_files:
                        logger.info(f"Reached maximum file limit ({max_files})")
                        break
                        
                    # Skip system files
                    if filename.startswith('.') or filename.startswith('__MACOSX'):
                        continue
                    
                    file_path = os.path.join(root, filename)
                    file_ext = os.path.splitext(filename)[1].lower()
                    
                    # Skip very large files
                    try:
                        if os.path.getsize(file_path) > 10 * 1024 * 1024:  # 10MB limit
                            logger.info(f"Skipping large file: {filename}")
                            continue
                    except:
                        continue
                    
                    logger.info(f"Processing ZIP file: {filename}")
                    
                    try:
                        with open(file_path, 'rb') as f:
                            content = f.read()

                        # Parse the file
                        text, metadata = DocumentParser.parse_document(content, file_ext)
                        if not text.strip() or len(text.strip()) < 10:
                            continue
                        
                        # Chunk the text
                        chunks, chunk_meta = SmartChunker.chunk_document(
                        text, 
                        [{'source': filename, 'type': f'zip_content_{file_ext}', **(metadata[0] if metadata else {})}], 
                        chunk_size=settings.CHUNK_SIZE_CHARS, 
                        overlap=settings.CHUNK_OVERLAP_CHARS
                        )
                        
                        # Add to vector store
                        if chunks:
                            embeddings = await pipeline._generate_embeddings(chunks)
                            vector_store.add(chunks, embeddings, chunk_meta)
                            processed_count += 1

                    except Exception as e:
                        logger.error(f"Failed to process '{filename}': {e}")
                        continue
                        
                if processed_count >= max_files:
                    break
                    
        except Exception as e:
            logger.error(f"ZIP processing failed: {e}")
            # Add error information to vector store
            fallback_chunks = [f"ZIP processing error: {str(e)[:200]}"]
            fallback_metadata = [{'source': zip_path, 'type': 'zip_error'}]
            embeddings = await pipeline._generate_embeddings(fallback_chunks)
            vector_store.add(fallback_chunks, embeddings, fallback_metadata)
            
        finally:
            # Clean up
            shutil.rmtree(temp_dir, ignore_errors=True)
    
    @staticmethod
    def parse_odt(content: bytes) -> Tuple[str, List[Dict]]:
        """Parse ODT files"""
        try:
            doc = load(io.BytesIO(content))
            text_parts = []
            
            for element in doc.getElementsByType(P):
                text = str(element)
                if text and text.strip():
                    text_parts.append(text)
            
            final_text = "\n\n".join(text_parts) if text_parts else "No text found in ODT"
            return final_text, [{'type': 'odt'}]
            
        except Exception as e:
            logger.error(f"ODT parsing failed: {e}")
            return f"Error parsing ODT: {str(e)}", [{'type': 'error'}]
    
    @staticmethod
    def parse_text(content: bytes) -> Tuple[str, List[Dict]]:
        """Fallback text extraction with encoding detection"""
        try:
            # Try multiple encodings
            for encoding in ['utf-8', 'latin-1', 'cp1252', 'ascii']:
                try:
                    text = content.decode(encoding, errors='ignore')
                    if text and text.strip():
                        return text, [{'type': 'text', 'encoding': encoding}]
                except:
                    continue
            
            return "Unable to decode text content", [{'type': 'error'}]
            
        except Exception as e:
            logger.error(f"Text parsing failed: {e}")
            return "Unable to extract text from file", [{'type': 'error'}]
    
    @staticmethod
    def _format_table(table: List[List]) -> str:
        """Format table data with better structure"""
        if not table:
            return ""
        
        formatted = []
        
        # Headers with proper formatting
        if len(table) > 0:
            headers = [str(cell) if cell else "" for cell in table[0]]
            formatted.append(" | ".join(headers))
            # Create separator line
            separator = " | ".join(["-" * max(3, len(h)) for h in headers])
            formatted.append(separator)
        
        # Data rows with null handling
        for row in table[1:]:
            if row and any(cell for cell in row):
                row_text = " | ".join([str(cell) if cell else "" for cell in row])
                formatted.append(row_text)
        
        return "\n".join(formatted)